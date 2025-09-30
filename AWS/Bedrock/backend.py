
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
#from pydantic import BaseSettings
from pydantic_settings import BaseSettings
from typing import List
import os, uuid, json, base64, mimetypes
import boto3
import chromadb
from chromadb.utils import embedding_functions
import PyPDF2
from PIL import Image
import io
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_bytes

# Recupera dinamicamente regione e account ID
session = boto3.session.Session()
aws_region = session.region_name or "eu-central-1"
account_id = session.client("sts").get_caller_identity()["Account"]

class Settings(BaseSettings):
    AWS_REGION: str = aws_region
    S3_BUCKET: str = "ragdemo-alnao-bucket"
    CHROMA_DIR: str = "./chroma"
    BEDROCK_EMBEDDING_MODEL: str = "amazon.titan-embed-text-v2:0"
    BEDROCK_CHAT_MODEL: str = "meta.llama3-2-3b-instruct-v1:0"  # adjust to what you enabled
    INFERENCE_PROFILE_ARN: str = f"arn:aws:bedrock:{aws_region}:{account_id}:inference-profile/eu.meta.llama3-2-3b-instruct-v1:0"
    USE_LOCAL_EMBEDDINGS: bool = True  # Fallback to local embeddings if Bedrock is not accessible
    BEDROCK_MODEL_ID_FOR_IMAGE_DESCRIPTION: str = "anthropic.claude-3-sonnet-20240229-v1:0"  # adjust to what you enabled

settings = Settings()

app = FastAPI(title="Pocket RAG")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, specify your domain instead of "*"
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Chroma (local, sqlite)
os.makedirs(settings.CHROMA_DIR, exist_ok=True)
client = chromadb.PersistentClient(path=settings.CHROMA_DIR)
# Use a different collection name for Bedrock embeddings to avoid dimension mismatch
collection = client.get_or_create_collection(name="docs-bedrock")

bedrock = boto3.client("bedrock-runtime", region_name=settings.AWS_REGION)


def bedrock_embed_batch(texts: List[str]) -> List[List[float]]:
    # Use Bedrock (original code)
    body = json.dumps({"inputText": texts})  # Titan v2 supports batch
    resp = bedrock.invoke_model(
        modelId=settings.BEDROCK_EMBEDDING_MODEL,
        body=body,
        contentType="application/json",
        accept="application/json",
    )
    payload = json.loads(resp["body"].read())
    return [item["embedding"] for item in payload["embedding"]]


def bedrock_embed(texts: List[str]) -> List[List[float]]:
    # Use Bedrock - process texts individually as Titan v2 expects single strings
    embeddings = []
    for text in texts:
        body = json.dumps({"inputText": text})
        resp = bedrock.invoke_model(
            modelId=settings.BEDROCK_EMBEDDING_MODEL,
            body=body,
            contentType="application/json",
            accept="application/json",
        )
        payload = json.loads(resp["body"].read())
        embeddings.append(payload["embedding"])
    return embeddings

def bedrock_chat(prompt: str) -> str:
    # Different body format for Llama 3.2 3B Instruct
    body = {
        "prompt": prompt,
        "temperature": 0.2,
        "top_p": 0.9
    }
    # Use inference profile ARN as modelId if available, otherwise use regular model ID
    model_id = settings.INFERENCE_PROFILE_ARN if settings.INFERENCE_PROFILE_ARN else settings.BEDROCK_CHAT_MODEL
    
    resp = bedrock.invoke_model(
        modelId=model_id,
        body=json.dumps(body),
        contentType="application/json",
        accept="application/json",
    )
    out = json.loads(resp["body"].read())
    # meta llama uses "output" or "content" depending on version; normalize:
    if "output" in out:  # some providers
        return out["output"]["message"]["content"][0]["text"]
    if "content" in out:
        return out["content"][0]["text"]
    return str(out)

def describe_image_with_bedrock(image_bytes: bytes, filename: str) -> str:
    """
    Describe an image using Claude 3 on Bedrock
    """
    try:
        # Convert image to base64
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
        
        # Determine media type
        media_type = "image/jpeg"
        if filename.lower().endswith('.png'):
            media_type = "image/png"
        elif filename.lower().endswith('.webp'):
            media_type = "image/webp"
        elif filename.lower().endswith('.gif'):
            media_type = "image/gif"
        
        body = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 1000,
            "messages": [{
                "role": "user",
                "content": [{
                    "type": "text",
                    "text": "Describe this image in detail. Extract any text you see and provide a comprehensive description of the visual content."
                }, {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": media_type,
                        "data": image_base64
                    }
                }]
            }]
        }
        
        resp = bedrock.invoke_model(
            modelId= settings.BEDROCK_MODEL_ID_FOR_IMAGE_DESCRIPTION,
            body=json.dumps(body),
            contentType="application/json",
            accept="application/json",
        )
        
        response_data = json.loads(resp["body"].read())
        return response_data["content"][0]["text"]
        
    except Exception as e:
        return f"Error describing image: {str(e)}"

def extract_text_from_file(file_content: bytes, filename: str, content_type: str) -> str:
    """
    Extract text from various file formats
    """
    filename_lower = filename.lower()
    
    try:
        # Text files
        if content_type.startswith('text/') or filename_lower.endswith(('.txt', '.md', '.csv', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.cpp', '.c', '.h')):
            return file_content.decode('utf-8', errors='ignore')
        
        # PDF files
        elif filename_lower.endswith('.pdf'):
            try:
                # Prima prova con PyPDF2
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                text = ""
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                #if text.strip():
                #    return text
                # Se non c'è testo, prova OCR su immagini delle pagine
                images = convert_from_bytes(file_content)
                ocr_text = text
                for img in images:
                    ocr_text += pytesseract.image_to_string(img, lang="ita+eng") + "\n"
                return ocr_text if ocr_text.strip() else "Error: Unable to extract text from PDF (no text or OCR failed)"
            except Exception as e:
                return f"Error: Unable to extract text from PDF: {str(e)}"
        
        # Image files - use Bedrock for description
        elif content_type.startswith('image/') or filename_lower.endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff')):
            return describe_image_with_bedrock(file_content, filename)
        
        # Word documents
        elif filename_lower.endswith('.docx'):
            try:
                doc_file = io.BytesIO(file_content)
                doc = Document(doc_file)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except:
                return "Error: Unable to extract text from Word document. Install python-docx: pip install python-docx"
        
        # Excel files
        elif filename_lower.endswith(('.xlsx', '.xls')):
            try:
                excel_file = io.BytesIO(file_content)
                workbook = openpyxl.load_workbook(excel_file)
                text = ""
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"Sheet: {sheet_name}\n"
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text += row_text + "\n"
                return text
            except:
                return "Error: Unable to extract text from Excel file. Install openpyxl: pip install openpyxl"
        
        # PowerPoint files
        elif filename_lower.endswith(('.pptx', '.ppt')):
            try:
                ppt_file = io.BytesIO(file_content)
                presentation = Presentation(ppt_file)
                text = ""
                for slide_num, slide in enumerate(presentation.slides, 1):
                    text += f"Slide {slide_num}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except:
                return "Error: Unable to extract text from PowerPoint file. Install python-pptx: pip install python-pptx"
        
        # Fallback: try to decode as text
        else:
            try:
                return file_content.decode('utf-8', errors='ignore')
            except:
                return f"Unsupported file type: {content_type}. Filename: {filename}"
                
    except Exception as e:
        return f"Error processing file {filename}: {str(e)}"

@app.get("/health") 
def health(): return {"ok": True}

@app.get("/supported-formats")
def supported_formats():
    return {
        "text_formats": [".txt", ".md", ".csv", ".json", ".xml", ".html", ".css", ".js", ".py", ".java", ".cpp", ".c", ".h"],
        "document_formats": [".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt"],
        "image_formats": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"],
        "notes": {
            "pdf": "Text extraction using PyPDF2",
            "images": "AI description using Claude 3 on Bedrock",
            "word": "Requires python-docx: pip install python-docx",
            "excel": "Requires openpyxl: pip install openpyxl", 
            "powerpoint": "Requires python-pptx: pip install python-pptx"
        }
    }

@app.post("/ingest")
async def ingest(file: UploadFile = File(...), doc_id: str = Form(default=None)):
    # Read file content
    file_content = await file.read()
    
    # Extract text based on file type
    text = extract_text_from_file(file_content, file.filename, file.content_type)
    
    # Check if text extraction was successful
    if text.startswith("Error:") or text.startswith("Unsupported"):
        return {"error": text, "filename": file.filename, "content_type": file.content_type}
    
    # naive splitter:
    chunks = [text[i:i+1200] for i in range(0, len(text), 1200)]
    
    # Filter out empty chunks
    chunks = [chunk.strip() for chunk in chunks if chunk.strip()]
    
    if not chunks:
        return {"error": "No text content found in file", "filename": file.filename}
    
    embeds = bedrock_embed(chunks)
    ids = [f"{doc_id or uuid.uuid4()}_{i}" for i in range(len(chunks))]
    
    # Add metadata with file type information
    metadata = [{"name": file.filename, "content_type": file.content_type, "chunk": i} for i in range(len(chunks))]
    
    collection.add(ids=ids, documents=chunks, embeddings=embeds, metadatas=metadata)
    
    # store raw in S3
    s3 = boto3.client("s3", region_name=settings.AWS_REGION)
    s3.put_object(Bucket=settings.S3_BUCKET, Key=f"uploads/{file.filename}", Body=file_content)
    
    return {
        "added": len(chunks), 
        "filename": file.filename, 
        "content_type": file.content_type,
        "extracted_text_length": len(text),
        "chunks_created": len(chunks)
    }

@app.post("/query")
async def query(q: str, k: int = 4):
    # Use Bedrock to embed the query
    query_embedding = bedrock_embed([q])[0]
    r = collection.query(query_embeddings=[query_embedding], n_results=k)
    ctx = "\n\n".join(r["documents"][0])
    prompt = f"Use the context to answer.\n\n# Context\n{ctx}\n\n# Question\n{q}\n\n# Answer:"
    answer = bedrock_chat(prompt)
    return {"answer": answer, "chunks": r["documents"][0], "ids": r["ids"][0]}

@app.get("/chroma/all")
def chroma_all():
    return collection.get()

@app.delete("/forget") #  /forget?filename=nomefile.ext
async def forget(filename: str):
    """
    Rimuove tutti i chunk, embedding e metadati associati a un file dalla collection ChromaDB.
    """
    # Trova tutti gli ID dei chunk associati a quel filename
    results = collection.get(where={"name": filename})
    ids_to_delete = results.get("ids", [])
    if not ids_to_delete:
        return {"status": "not_found", "message": f"Nessun chunk trovato per il file {filename}"}
    # Elimina i chunk dalla collection
    collection.delete(ids=ids_to_delete)
    s3 = boto3.client("s3", region_name=settings.AWS_REGION)
    s3.delete_object(Bucket=settings.S3_BUCKET, Key=f"uploads/{filename}")
    return {"status": "ok", "deleted_chunks": len(ids_to_delete), "filename": filename}

@app.post("/generate")
async def generate(prompt: str):
    return {"answer": bedrock_chat(prompt)}